from langflow.custom import Component
from langflow.io import MessageTextInput, StrInput, SecretStrInput, DropdownInput, Output
from langflow.schema.message import Message


class OCRProcessorNode(Component):
    display_name = "OCR Processor (AI Vision)"
    description = "Extracts text from an image or PDF using an AI vision model (GPT-4o, Claude, etc.)."
    icon = "scan"

    inputs = [
        MessageTextInput(name="file_path", display_name="File Path(s)", info="Path(s) to image or PDF file(s). Multiple paths separated by newline."),
        DropdownInput(
            name="provider",
            display_name="LLM Provider",
            options=["openai", "anthropic"],
            value="openai",
        ),
        StrInput(
            name="base_url",
            display_name="Base URL (optional)",
            info="Custom API base URL (e.g. for vLLM). Leave empty for default.",
            value="",
            advanced=True,
        ),
        StrInput(name="model_name", display_name="Model Name", value="gpt-4o"),
        SecretStrInput(name="api_key", display_name="API Key", value="not-needed", info="For vLLM/Ollama enter any value, e.g. 'not-needed'."),
        StrInput(
            name="prompt",
            display_name="OCR Prompt",
            value="Extract ALL text from this document image. Return only the raw text, preserving the structure. Do not add any commentary.",
            advanced=True,
        ),
    ]

    outputs = [
        Output(name="extracted_text", display_name="Extracted Text", method="run_ocr"),
    ]

    def _image_to_base64(self, image_path: str) -> tuple[str, str]:
        import base64
        import os

        ext = os.path.splitext(image_path)[1].lower()
        media_type_map = {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".webp": "image/webp",
            ".bmp": "image/bmp",
        }
        media_type = media_type_map.get(ext, "image/png")

        with open(image_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("utf-8")

        return b64, media_type

    def _pdf_to_images(self, pdf_path: str) -> list[str]:
        import tempfile

        tmp_dir = tempfile.mkdtemp(prefix="langflow_ocr_")
        image_paths = []

        try:
            import fitz  # PyMuPDF

            doc = fitz.open(pdf_path)
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                img_path = f"{tmp_dir}/page_{i}.png"
                pix.save(img_path)
                image_paths.append(img_path)
            doc.close()
        except ImportError:
            try:
                from pdf2image import convert_from_path

                images = convert_from_path(pdf_path, dpi=150)
                for i, img in enumerate(images):
                    img_path = f"{tmp_dir}/page_{i}.png"
                    img.save(img_path, "PNG")
                    image_paths.append(img_path)
            except ImportError:
                image_paths.append(pdf_path)

        return image_paths

    def _ocr_openai(self, image_paths: list[str]) -> str:
        import openai

        api_key = self.api_key if self.api_key else "not-needed"
        kwargs = {"api_key": api_key}
        if self.base_url:
            base_url = self.base_url.strip().rstrip("/")
            if not base_url.endswith("/v1"):
                base_url += "/v1"
            kwargs["base_url"] = base_url

        client = openai.OpenAI(**kwargs)

        # Process each page separately to avoid exceeding token limits
        results = []
        for img_path in image_paths:
            b64, media_type = self._image_to_base64(img_path)
            content = [
                {"type": "text", "text": self.prompt},
                {"type": "image_url", "image_url": {"url": f"data:{media_type};base64,{b64}", "detail": "high"}},
            ]
            response = client.chat.completions.create(
                model=self.model_name,
                max_tokens=2000,
                messages=[{"role": "user", "content": content}],
            )
            results.append(response.choices[0].message.content)

        return "\n\n--- Page Break ---\n\n".join(results)

    def _ocr_anthropic(self, image_paths: list[str]) -> str:
        import anthropic

        client = anthropic.Anthropic(api_key=self.api_key)

        results = []
        for img_path in image_paths:
            b64, media_type = self._image_to_base64(img_path)
            content = [
                {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": b64}},
                {"type": "text", "text": self.prompt},
            ]
            response = client.messages.create(
                model=self.model_name,
                max_tokens=2000,
                messages=[{"role": "user", "content": content}],
            )
            results.append(response.content[0].text)

        return "\n\n--- Page Break ---\n\n".join(results)

    def _extract_text_file(self, file_path: str) -> str:
        """Extract text from non-image files (DOCX, XLSX, CSV, TXT, etc.)."""
        import os
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()

        if ext == ".csv":
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()

        if ext == ".docx":
            try:
                from docx import Document
                doc = Document(file_path)
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            except ImportError:
                print("[OCR] python-docx not installed, falling back to OCR for .docx")
                return ""

        if ext == ".xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                parts = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    parts.append(f"[Arkusz: {sheet}]")
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        if any(cells):
                            parts.append(" | ".join(cells))
                wb.close()
                return "\n".join(parts)
            except ImportError:
                print("[OCR] openpyxl not installed, falling back to OCR for .xlsx")
                return ""

        if ext == ".xls":
            try:
                import xlrd
                wb = xlrd.open_workbook(file_path)
                parts = []
                for sheet in wb.sheet_names():
                    ws = wb.sheet_by_name(sheet)
                    parts.append(f"[Arkusz: {sheet}]")
                    for row_idx in range(ws.nrows):
                        cells = [str(ws.cell_value(row_idx, col)) for col in range(ws.ncols)]
                        if any(cells):
                            parts.append(" | ".join(cells))
                return "\n".join(parts)
            except ImportError:
                print("[OCR] xlrd not installed, falling back to OCR for .xls")
                return ""

        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                parts = []
                for i, slide in enumerate(prs.slides, 1):
                    parts.append(f"[Slajd {i}]")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    parts.append(para.text)
                return "\n".join(parts)
            except ImportError:
                print("[OCR] python-pptx not installed, falling back to OCR for .pptx")
                return ""

        return ""

    def _ocr_single_file(self, file_path: str) -> str:
        import os

        if not file_path or not os.path.exists(file_path):
            return ""

        ext = os.path.splitext(file_path)[1].lower()

        # Try text extraction first for supported formats
        text_formats = {".txt", ".csv", ".docx", ".xlsx", ".xls", ".pptx"}
        if ext in text_formats:
            text = self._extract_text_file(file_path)
            if text:
                return text

        if ext == ".pdf":
            image_paths = self._pdf_to_images(file_path)
        else:
            image_paths = [file_path]

        if self.provider == "anthropic":
            return self._ocr_anthropic(image_paths)
        else:
            return self._ocr_openai(image_paths)

    def run_ocr(self) -> Message:
        raw = self.file_path
        if isinstance(raw, Message):
            raw = raw.text

        file_paths = [p.strip() for p in raw.strip().split("\n") if p.strip()]

        results = []
        for fp in file_paths:
            text = self._ocr_single_file(fp)
            if text:
                results.append(f"[{fp}]\n{text}")

        return Message(text="\n\n===\n\n".join(results))
